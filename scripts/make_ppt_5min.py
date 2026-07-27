# -*- coding: utf-8 -*-
"""5 分钟验收答辩 PPT：11 页，全部 8 张图嵌入。深蓝主题 + 金色点缀。"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree as _etree

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGS = os.path.join(ROOT, "docs", "figs")
sys.stdout.reconfigure(errors="replace")

def OxmlElement(nsptag_str):
    from pptx.oxml.ns import NamespacePrefixedTag
    return _etree.Element(NamespacePrefixedTag(nsptag_str).clark_name)

# ═══════ 设计系统 ═══════
PRIMARY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT  = RGBColor(0xC9, 0xA9, 0x6E)
BLUE   = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
GRAY75 = RGBColor(0x75, 0x75, 0x75)
GRAY88 = RGBColor(0x88, 0x88, 0x88)
GRAY95 = RGBColor(0x95, 0x95, 0x95)
GRAYF5 = RGBColor(0xF5, 0xF5, 0xF5)
GRAYE0 = RGBColor(0xE0, 0xE0, 0xE0)

W = Inches(13.333); H = Inches(7.5)

def bg(slide, color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None):
    s=slide.shapes.add_shape(1,l,t,w,h); s.line.fill.background()
    if fill: s.fill.solid(); s.fill.fore_color.rgb=fill
    return s

def left_bar(slide, l, t, h, c=ACCENT): rect(slide,l,t,Inches(0.06),h,c)
def bottom_bar(slide): rect(slide,Inches(0),H-Inches(0.06),W,Inches(0.06),PRIMARY)

def pno(slide, n, t=11):
    tb=slide.shapes.add_textbox(Inches(11.8),H-Inches(0.55),Inches(1.2),Inches(0.4))
    p=tb.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT
    r=p.add_run(); r.text=f"{n}/{t}"; r.font.size=Pt(8); r.font.color.rgb=GRAY88; r.font.name="Arial"

def tag(slide, txt):
    tb=slide.shapes.add_textbox(Inches(0.55),H-Inches(0.55),Inches(3),Inches(0.4))
    r=tb.text_frame.paragraphs[0].add_run(); r.text=txt; r.font.size=Pt(8); r.font.color.rgb=GRAY95; r.font.name="Microsoft YaHei"

def title_text(slide, txt, y=Inches(0.45)):
    left_bar(slide, Inches(0.55), y+Inches(0.08), Inches(0.48))
    tb=slide.shapes.add_textbox(Inches(0.8), y, Inches(11.5), Inches(0.6))
    r=tb.text_frame.paragraphs[0].add_run(); r.text=txt
    r.font.size=Pt(24); r.font.color.rgb=PRIMARY; r.font.bold=True; r.font.name="Microsoft YaHei"

def bullets(slide, items, x, y, w, h, fs=11.5, sp=6):
    tb=slide.shapes.add_textbox(x, y, w, h); tf=tb.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(sp)
        dot_r=p.add_run(); dot_r.text="  ●  "; dot_r.font.size=Pt(9); dot_r.font.color.rgb=ACCENT
        if isinstance(item,tuple):
            r=p.add_run(); r.text=item[0]; r.font.size=Pt(fs); r.font.color.rgb=item[1]; r.font.name="Microsoft YaHei"
        else:
            r=p.add_run(); r.text=item; r.font.size=Pt(fs); r.font.color.rgb=DARK; r.font.name="Microsoft YaHei"

def fig(slide, name, l, t, w, cap=None):
    fp=os.path.join(FIGS,name)
    if os.path.exists(fp): slide.shapes.add_picture(fp,l,t,width=w)
    if cap:
        tb=slide.shapes.add_textbox(l,t+Inches(3.4),w,Inches(0.3))
        p=tb.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=cap; r.font.size=Pt(8); r.font.color.rgb=GRAY75; r.font.name="Microsoft YaHei"

def table(slide, rows, l, t, cw, fs=8):
    nr,nc=len(rows),len(rows[0]); tw=sum(cw)
    tbl=slide.shapes.add_table(nr,nc,l,t,tw,Inches(0.30*nr))
    tab=tbl.table
    for j,w_ in enumerate(cw): tab.columns[j].width=w_
    for i,row in enumerate(rows):
        for j,val in enumerate(row):
            cell=tab.cell(i,j); cell.text=""
            p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
            r=p.add_run(); r.text=str(val); r.font.size=Pt(fs); r.font.name="Arial"
            r.font.bold=(i==0); r.font.color.rgb=WHITE if i==0 else DARK
            cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            if i==0: cell.fill.solid(); cell.fill.fore_color.rgb=PRIMARY
            else: cell.fill.solid(); cell.fill.fore_color.rgb=GRAYF5 if i%2==0 else WHITE

def std_slide(slide, title, items, fig_name=None, section="", n=1, fs=11.5):
    bg(slide, GRAYF5); bottom_bar(slide); pno(slide,n); tag(slide,section)
    title_text(slide,title)
    has_fig=fig_name is not None
    bw=Inches(7.2) if has_fig else Inches(12.2)
    cy,ch=Inches(1.3), Inches(5.5)
    rect(slide, Inches(0.55), cy, bw, ch, WHITE)
    bullets(slide, items, Inches(0.85), Inches(cy.inches+0.15), Inches(bw.inches-0.6), Inches(ch.inches-0.3), fs=fs)
    if has_fig:
        fi=Inches(7.7) if fig_name.startswith("fig5") else Inches(8.0)
        fw=Inches(5.2) if fig_name in ("fig5_ablation.png","fig8_dist.png","fig6_progression.png","fig4_scatter.png") else Inches(4.8)
        if fig_name in ("fig1_pipeline.png",): fi,fw=Inches(7.7), Inches(5.0)
        fig(slide, fig_name, fi, Inches(1.55), fw)


# ═══════ 11 页 ═══════
def build():
    prs = Presentation()
    prs.slide_width=W; prs.slide_height=H; bl=prs.slide_layouts[6]

    # P1 封面 (10s)
    s=prs.slides.add_slide(bl); bg(s,PRIMARY)
    rect(s,Inches(8.5),Inches(0),Inches(5),Inches(7.5),RGBColor(0x21,0x36,0x5A))
    rect(s,Inches(4.2),Inches(2.8),Inches(1.8),Inches(0.01),ACCENT)
    tb=s.shapes.add_textbox(Inches(1.2),Inches(1.6),Inches(6.8),Inches(2.2)); tf=tb.text_frame; tf.word_wrap=True
    r=tf.paragraphs[0].add_run(); r.text="基于开源多模态大模型的"; r.font.size=Pt(30); r.font.color.rgb=WHITE; r.font.name="Microsoft YaHei"
    r=tf.add_paragraph().add_run(); r.text="图像质量评估 Agent 框架"; r.font.size=Pt(36); r.font.color.rgb=WHITE; r.font.bold=True; r.font.name="Microsoft YaHei"
    tf.add_paragraph().space_after=Pt(15)
    r=tf.add_paragraph().add_run(); r.text="Unified Agent Framework for Blind IQA"; r.font.size=Pt(15); r.font.color.rgb=ACCENT; r.font.name="Arial"
    tf.add_paragraph().space_after=Pt(25)
    r=tf.add_paragraph().add_run(); r.text="张丞  ·  北京邮电大学"; r.font.size=Pt(18); r.font.color.rgb=WHITE; r.font.name="Microsoft YaHei"
    r=tf.add_paragraph().add_run(); r.text="指导教师：何老师    验收汇报  ·  2026.07.29"; r.font.size=Pt(13); r.font.color.rgb=GRAY88; r.font.name="Microsoft YaHei"

    # P2 任务与约束 (25s)
    std_slide(prs.slides.add_slide(bl), "任务与硬约束", [
        "任务：构建通用 IQA Agent 框架 —— Backbone 零训练，仅通过 Skill（Prompt）挖掘大模型质量评估能力",
        "主干：Qwen3-VL-32b-instruct（DashScope API）· 温度 = 0 · 全流程 SHA256 磁盘缓存可复现",
        "评测：KonIQ-10k Val 2,015 张（1–5 整数 MOS）· SPAQ Test 1,125 张（0–10 连续 MOS）",
        "合规红线（任务书 §4.1–§4.5）：",
        "  → 数据集名 / Image ID / MOS / 失真类型等级 / 划分信息 — 一律不进入 Router/LLM/Skill 线上输入",
        "  → MOS 仅用于离线评测与误差分析 —— 全程每轮末各读一次，均在预注册文档冻结后执行",
        "  → 仅 Router/Decision 层可训练（3 项职能：规则选择 · 冲突裁决 · 解释生成）",
        ("监督信号全部来自训练集像素自衍生（合成失真阶梯 + 两两比较 BT 锦标赛）— 零 MOS 接触", GRAY75),
    ], section="背景", n=2)

    # P3 框架总览 + 图1 (35s)
    std_slide(prs.slides.add_slide(bl), "统一框架总览：一张图，三路并行", [
        "三路并行输入：",
        "  ① 技能专家并行评分（S-TECH / S-GLOBAL / S-CONTENT）",
        "  ② 像素统计特征 F（7 维 OpenCV：清晰度/噪声/色彩/亮度/分辨率/宽高比/分歧度）",
        "  ③ 裸问释义投票（4 条同义问法取均值 —— 承担零点校准）",
        "门控矩阵 W(3×7)：BT 锦标赛离线训练 → 冻结后纯矩阵乘法推理 → 输出动态话语权 g = softmax(W·F)",
        "融合分 s_fus = g · s（三专家分加权内积）；投票分 s_vote = mean（4 条释义 + 裸问原句）",
        "最终分 = α · s_fus + (1−α) · s_vote  —— KonIQ α=0.6 / SPAQ α=0.3（小店扫描）",
        "R1→R2→R3 搭建三路，R4→R5 预注册验证，R6 统一收敛 —— 两数据集仅参数不同",
    ], "fig1_pipeline.png", "框架", 3, fs=11)

    # P4 BT 锦标赛 + 门控训练（两张图左右各一）(45s)
    s=prs.slides.add_slide(bl)
    bg(s, GRAYF5); bottom_bar(s); pno(s,4); tag(s,"方法")
    title_text(s, "核心方法：BT 锦标赛自监督 + 门控矩阵训练")
    # 左栏 BT
    bullets(s, [
        "监督信号难题：路由需训练信号，但人群分被 §4.4 禁止 → 答案：两两比较锦标赛",
        "科学依据：VLM 绝对评分噪声 ‖ 两两比较能力（SPAQ梯子：绝对值全FAIL ‖ 比较0.98/0.88）— 量级差距",
        "组织：2,500 张 KonIQ Train 图 → 32,389 场两两比较 → Bradley-Terry 凸优化 → 每图一标量强度分",
        "本质：把大模型内部自发的质量排序能力蒸馏成训练信号 —— 零人工标签、不依赖任何 MOS",
    ], Inches(0.7), Inches(1.4), Inches(5.8), Inches(2.8), fs=11)
    # 左下图: loss曲线
    fig(s, "fig2_loss.png", Inches(0.55), Inches(4.05), Inches(5.9))
    # 右栏 门控
    bullets(s, [
        "门控训练：W(3×7) 从零初始化 → 每步 128 对图 → 解析梯度 SGD → 800 步 CPU 数秒收敛",
        "损失 = −log σ(s_fus_A − s_fus_B)：最小化成对铰链损失，优化方向 = 满足 BT 偏好顺序",
    ], Inches(7.0), Inches(1.4), Inches(5.8), Inches(1.5), fs=11)
    # 右下图: W矩阵
    fig(s, "fig3_w.png", Inches(7.0), Inches(3.3), Inches(5.8))
    bullets(s, [
        ("训练后 W 高度可读：分歧度列权重绝对值最大（TECH +0.89 / CONTENT −0.84）", ACCENT),
        ("→ 专家意见不一致时，门控把话语权交给技术专家；logpix/aspect权重自动归零（尺寸恒定无信息）", GRAY75),
        "留出 2,000 对检验：动态融合一致率 80.9% vs 等权 78.0%（+2.9 pp）",
    ], Inches(7.0), Inches(5.6), Inches(5.8), Inches(1.6), fs=10.5)

    # P5 结果主表 + 趋势图 (40s)
    s=prs.slides.add_slide(bl)
    bg(s,GRAYF5); bottom_bar(s); pno(s,5); tag(s,"结果")
    title_text(s, "三轮综合主表与逐臂趋势")
    hdr=["轮","臂","KonIQ SRCC","KonIQ MAE","SPAQ SRCC","SPAQ MAE"]
    rows=[hdr,
          ["一","R1-bare（裸问基线）","0.577","0.463","0.881","0.893"],
          ["一","R1-rich（+完整细则）","0.633","0.688","0.867","1.042"],
          ["一","R2.5（+动态分诊）","0.619","0.533","0.860","1.014"],
          ["二","R4（BT监督路由）","0.668","0.554","0.885","0.908"],
          ["三","R6（统一框架）","0.734","0.491","0.893","0.945"]]
    table(s, rows, Inches(0.4), Inches(1.5),
          [Inches(0.5),Inches(2.6),Inches(1.35),Inches(1.3),Inches(1.35),Inches(1.3)], fs=9)
    fig(s, "fig6_progression.png", Inches(7.5), Inches(1.5), Inches(5.4))
    bullets(s, [
        "R6 相对 R1-bare 基准：KonIQ SRCC +0.157、MAE −0.028；SPAQ SRCC +0.012",
        "R2.5→R4 +0.049：同分布 BT 信号替换旧梯子（同一机制、不同信号、收益近一个量级）",
        "R4→R6 +0.065：逐图动态门控（+0.026一致率）+ 裸问释义投票零点校准（+0.045 τ）",
        ("CKE 经验规则（R2.5→R3）SRCC仅 +0.002 —— 自洽≠对齐 —— 催生锚相关性定律", GRAY75),
        "对比：Tool-IQA 同族 8B 零样本 SRCC 0.729 —— 本文 32B 以 0.734 越过此锚点",
    ], Inches(0.55), Inches(4.1), Inches(12.2), Inches(2.6), fs=11)

    # P6 效果图 + 迁移检验 (25s)
    s=prs.slides.add_slide(bl)
    bg(s,GRAYF5); bottom_bar(s); pno(s,6); tag(s,"检验")
    title_text(s, "散点密度 · 门控跨模型迁移（0 新增调用）")
    fig(s, "fig4_scatter.png", Inches(0.45), Inches(1.5), Inches(6.3))
    fig(s, "fig7_transfer.png", Inches(6.95), Inches(1.5), Inches(5.8))
    bullets(s, [
        "左：R6 预测分 vs 人群分散点密度 —— 两数据集主趋势贴合对角线，无系统性偏移",
        "右：32B 训练的门控矩阵 W 原样套到 8B 专家分 —— 零新增 API、纯矩阵乘法",
        "8B 等权门控 0.759 → 迁移 W 后 0.776，反超 32B 自身（0.734）",
        ("→ 门控学到的是\"图像属性 → 专家可信度\"的通用映射，与 Backbone 规模解耦", ACCENT),
    ], Inches(0.55), Inches(5.2), Inches(12.2), Inches(1.8), fs=11)

    # P7 困难与发现：消融 + 分布 (50s)
    s=prs.slides.add_slide(bl)
    bg(s,GRAYF5); bottom_bar(s); pno(s,7); tag(s,"发现")
    title_text(s, "困难与发现：提示词约束的保序-校零分离")
    fig(s, "fig5_ablation.png", Inches(0.45), Inches(1.5), Inches(6.3))
    # 右上消融表
    hdr2=["臂","KonIQ MAE","SPAQ MAE","说明"]
    rows2=[hdr2,
           ["R1-bare(纯裸问)","0.454","0.819","零修饰先验"],
           ["R1-anchor v2","0.507","1.034","+人设+锚点带"],
           ["R1-anchor v3","0.645","1.062","+维度/清单/程序"],
           ["R1-rich","0.688","1.042","完整专家人设"]]
    table(s, rows2, Inches(7.2), Inches(1.5),
          [Inches(1.5),Inches(1.1),Inches(1.1),Inches(1.8)], fs=8.5)
    # 右下分布图
    fig(s, "fig8_dist.png", Inches(7.2), Inches(3.4), Inches(5.8))
    bullets(s, [
        "核心发现：排序（SRCC）与绝对校准（MAE）是可分离的两种能力",
        "一级级加约束 → SRCC 高位维持，MAE 单调恶化 —— 锚点带入排序知识同时框死输出分布的零点",
        "8B vs 32B 对比：纯裸问 MAE 0.579 vs 0.463 → 对人群趋势的把握来自模型能力本身，非 prompt 工程",
        ("这不是缺陷：R6 中裸问投票独立承担零点校准 —— \"锚点保序、裸问校零\"的互补结构", ACCENT),
    ], Inches(0.55), Inches(5.05), Inches(12.2), Inches(2.0), fs=10.5)
    # 右侧分布图说明
    bullets(s, [
        ("镜像证据：完整专家输出在锚点带中心形成尖峰（分布图），裸问更贴近人群分形态", GRAY75),
    ], Inches(7.2), Inches(6.58), Inches(5.8), Inches(0.6), fs=9)

    # P8 锚点双重角色 + 其他发现 (25s)
    std_slide(prs.slides.add_slide(bl), "困难与发现（续）：锚点双重角色 · CKE 零收益 · MAE 零点偏置", [
        "锚点双重角色（R1-anchor 探针 + R6-offanchor 消融）：",
        "  → 单问场景（R1-bare→R1-anchor）：锚点单独有害 —— MAE 0.463→0.544",
        "  → 多专家融合（R6→R6-offanchor）：锚点联合有益 —— 砍锚后 SPAQ MAE 0.945→1.208 崩盘（跨专家分不在同一天平）",
        ("  → 同一组件在不同架构深度扮演相反角色 —— 本项目对\"零训练下如何设计评估架构\"的精确定量回答", ACCENT),
        "CKE 自进化规则库（无效组件）：内部指标全改善（阶梯单调性 +0.019 / B-C 分歧 −0.07），外部 SRCC 仅 +0.001",
        "  → ¥69 换来零收益 —— 催生锚相关性定律（第一块基石）：监督信号与目标的相关性决定自监督上限",
        "MAE 零点偏置（结构性边界）：模型先验内容依赖零点摆幅 ~19pp ‖ 标注人群 ~6pp，在 §4.5 约束下原则性不可达",
    ], section="发现", n=8, fs=10.5)

    # P9 锚相关性定律闭合 (15s)
    s=prs.slides.add_slide(bl)
    bg(s,PRIMARY); rect(s,Inches(0),Inches(0),Inches(0.06),H,ACCENT); rect(s,Inches(0),H-Inches(0.06),W,Inches(0.06),ACCENT); pno(s,9)
    tb=s.shapes.add_textbox(Inches(1.0),Inches(0.6),Inches(11.3),Inches(6.3)); tf=tb.text_frame; tf.word_wrap=True
    r=tf.paragraphs[0].add_run(); r.text="锚相关性定律：三轮完整闭合"; r.font.size=Pt(28); r.font.color.rgb=WHITE; r.font.bold=True; r.font.name="Microsoft YaHei"
    tf.add_paragraph().space_after=Pt(20)
    items=[("第一轮 · 发现","CKE +0.001 零收益证伪\"自洽驱动自监督\"","同源多专家 +0.007 证伪\"旧信号上重加权\"","→ 确立\"信号相关性决定上限\"的定性命题"),
           ("第二轮 · 验证","同一融合机制：旧梯子信号 +0.007 ‖ BT信号 +0.062","同一代码、不同信号、收益差近一个量级","→ 定性命题升级为可操作的定量定律"),
           ("第三轮 · 放大","扩容锦标赛 32,000 场对决 · 引入逐图动态门控(21参数) · 增设裸问投票校零","KonIQ SRCC 0.577→0.734  SPAQ 0.881→0.893","跨模型迁移(32B W→8B,SRCC 0.776)验证定律的去模型依赖性")]
    for title, l1, l2, l3 in items:
        r=tf.add_paragraph().add_run(); r.text=title; r.font.size=Pt(18); r.font.color.rgb=ACCENT; r.font.bold=True; r.font.name="Microsoft YaHei"
        for l in [l1,l2,l3]:
            r=tf.add_paragraph().add_run(); r.text=f"    {l}"; r.font.size=Pt(13); r.font.color.rgb=WHITE; r.font.name="Microsoft YaHei"
        tf.add_paragraph().space_after=Pt(4)

    # P10 结论 (15s)
    s=prs.slides.add_slide(bl); bg(s,PRIMARY)
    rect(s,Inches(0),Inches(0),Inches(0.06),H,ACCENT); rect(s,Inches(0),H-Inches(0.06),W,Inches(0.06),ACCENT); pno(s,10)
    tb=s.shapes.add_textbox(Inches(1.2),Inches(1.0),Inches(11.0),Inches(4.5)); tf=tb.text_frame; tf.word_wrap=True
    r=tf.paragraphs[0].add_run(); r.text="结论"; r.font.size=Pt(32); r.font.color.rgb=WHITE; r.font.bold=True; r.font.name="Microsoft YaHei"
    tf.add_paragraph().space_after=Pt(20)
    for n,txt in [("01","统一框架：两数据集 × 两模型规模，同一骨架，主干零训练，全程温度=0 缓存可复现"),
                   ("02","BT 锦标赛自监督 → 21 参数可解释门控矩阵 → 有效（留出验证）· 可迁移（跨模型反超原主）"),
                   ("03","三项核心发现：提示词约束的保序-校零分离 · 锚点的双重角色 · 检索式经验规则的零收益"),
                   ("04","锚相关性定律：三轮完整闭合（发现→验证→放大——预注册假设检验保证可归因性）")]:
        r=tf.add_paragraph().add_run(); r.text=n; r.font.size=Pt(28); r.font.color.rgb=ACCENT; r.font.bold=True; r.font.name="Arial"
        r=tf.add_paragraph().add_run(); r.text=txt; r.font.size=Pt(15); r.font.color.rgb=WHITE; r.font.name="Microsoft YaHei"
        tf.add_paragraph().space_after=Pt(10)

    # P11 结尾 (5s)
    s=prs.slides.add_slide(bl); bg(s,PRIMARY)
    rect(s,Inches(0),Inches(0),Inches(0.06),H,ACCENT); rect(s,Inches(0),H-Inches(0.06),W,Inches(0.06),ACCENT); pno(s,11)
    tb=s.shapes.add_textbox(Inches(1.2),Inches(1.8),Inches(11.0),Inches(4.0)); tf=tb.text_frame; tf.word_wrap=True
    r=tf.paragraphs[0].add_run(); r.text="谢谢各位老师，请批评指正"; r.font.size=Pt(36); r.font.color.rgb=WHITE; r.font.bold=True; r.font.name="Microsoft YaHei"
    tf.add_paragraph().space_after=Pt(20)
    r=tf.add_paragraph().add_run(); r.text="张丞  ·  北京邮电大学  ·  2026.07.29"; r.font.size=Pt(18); r.font.color.rgb=ACCENT; r.font.name="Microsoft YaHei"
    r=tf.add_paragraph().add_run(); r.text="材料与代码见附件  ·  docs/findings.md（F-001~F-025 全部证据链）"; r.font.size=Pt(12); r.font.color.rgb=GRAY88; r.font.name="Microsoft YaHei"

    path=os.path.join(ROOT,"docs","汇报_5min.pptx")
    prs.save(path); print("->",path)

if __name__=="__main__":
    build()
