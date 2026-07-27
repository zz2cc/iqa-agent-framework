# -*- coding: utf-8 -*-
"""学术答辩 PPT v3：33 页详版，每章 3-5 页。深蓝主题 + 金色点缀 + 卡片式布局。
设计原则融合 ppt-master（序列执行、SVG一级设计契约→此处PPTX一级、严格门控路由）。"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree as _etree

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGS = os.path.join(ROOT, "docs", "figs")
sys.stdout.reconfigure(errors="replace")

def OxmlElement(nsptag_str):
    from pptx.oxml.ns import NamespacePrefixedTag
    return _etree.Element(NamespacePrefixedTag(nsptag_str).clark_name)

# ═══════ 设计系统（ppt-master规范）═══════
PRIMARY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT  = RGBColor(0xC9, 0xA9, 0x6E)
BLUE_A  = RGBColor(0x2A, 0x78, 0xD6)
ORANGE  = RGBColor(0xEB, 0x68, 0x34)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x1A, 0x1A)
GRAY75  = RGBColor(0x75, 0x75, 0x75)
GRAY88  = RGBColor(0x88, 0x88, 0x88)
GRAY95  = RGBColor(0x95, 0x95, 0x95)
GRAYE0  = RGBColor(0xE0, 0xE0, 0xE0)
GRAYF0  = RGBColor(0xF0, 0xF0, 0xF0)
GRAYF5  = RGBColor(0xF5, 0xF5, 0xF5)
RED_A   = RGBColor(0xD0, 0x3B, 0x3B)
GREEN_A = RGBColor(0x0C, 0xA3, 0x0C)

W = Inches(13.333); H = Inches(7.5)

def _set_slide_bg(slide, color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, lw=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color; shape.line.width = lw or Pt(1)
    return shape

def add_left_accent(slide, left, top, height, color=ACCENT, width=Inches(0.06)):
    return add_rect(slide, left, top, width, height, fill_color=color)

def add_bottom_bar(slide, color=PRIMARY):
    add_rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), fill_color=color)

def add_page_no(slide, num, total=33):
    tb = slide.shapes.add_textbox(Inches(11.8), H - Inches(0.55), Inches(1.2), Inches(0.4))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = f"{num} / {total}"
    r.font.size = Pt(8); r.font.color.rgb = GRAY88; r.font.name = "Arial"

def add_section_tag(slide, tag, x=Inches(0.55), y=H - Inches(0.55)):
    tb = slide.shapes.add_textbox(x, y, Inches(3), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = tag; r.font.size = Pt(8); r.font.color.rgb = GRAY95
    r.font.name = "Microsoft YaHei"

def add_title(slide, text, y=Inches(0.55)):
    add_left_accent(slide, Inches(0.55), y + Inches(0.08), Inches(0.52))
    tb = slide.shapes.add_textbox(Inches(0.85), y, Inches(11.5), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(26); r.font.color.rgb = PRIMARY; r.font.bold = True
    r.font.name = "Microsoft YaHei"

def add_subtitle(slide, text, x=Inches(0.85), y=Inches(1.25)):
    tb = slide.shapes.add_textbox(x, y, Inches(11.5), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(14); r.font.color.rgb = GRAY75; r.font.name = "Microsoft YaHei"

def add_bullets(slide, items, x=Inches(0.85), y=Inches(1.45), w=Inches(11.8), h=Inches(5.2), fs=13, spacing=8, color=DARK):
    """添加带金色圆点的要点列表。每项可以是 str 或 (str, RGBColor)。"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        r = p.add_run(); r.text = "  ●  "; r.font.size = Pt(10); r.font.color.rgb = ACCENT
        if isinstance(item, tuple):
            r = p.add_run(); r.text = item[0]; r.font.size = Pt(fs)
            r.font.color.rgb = item[1]; r.font.name = "Microsoft YaHei"
        else:
            r = p.add_run(); r.text = item; r.font.size = Pt(fs)
            r.font.color.rgb = color; r.font.name = "Microsoft YaHei"

def add_kpi_box(slide, left, top, width, label, value, delta=None, color=DARK):
    add_rect(slide, left, top, width, Inches(0.85), fill_color=WHITE)
    tb = slide.shapes.add_textbox(Inches(left.inches + 0.12), Inches(top.inches + 0.08), Inches(width.inches - 0.24), Inches(0.75))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(value); r.font.size = Pt(22); r.font.color.rgb = color; r.font.bold = True
    r.font.name = "Arial"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label; r2.font.size = Pt(9); r2.font.color.rgb = GRAY75
    r2.font.name = "Microsoft YaHei"
    if delta:
        p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run(); r3.text = delta; r3.font.size = Pt(8); r3.font.color.rgb = GREEN_A if "+" in delta else RED_A
        r3.font.name = "Arial"

def make_table(slide, rows, left, top, col_widths, header_color=PRIMARY, fs=8.5):
    n_rows, n_cols = len(rows), len(rows[0])
    total_w = sum(col_widths)
    row_h = Inches(0.32) if n_rows > 6 else Inches(0.35)
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, row_h * n_rows)
    table = tbl.table
    for j, w in enumerate(col_widths):
        table.columns[j].width = w
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(fs); r.font.name = "Arial"; r.font.bold = (i == 0)
            r.font.color.rgb = WHITE if i == 0 else DARK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = GRAYF5 if i % 2 == 0 else WHITE
            tcPr = cell._tc.get_or_add_tcPr()
            borders = _etree.Element("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tcBorders")
            for edge in ["top","bottom","left","right"]:
                el = _etree.SubElement(borders, f"{{http://schemas.openxmlformats.org/wordprocessingml/2006/main}}{edge}")
                el.set(qn("w:val"), "single"); el.set(qn("w:sz"), "4"); el.set(qn("w:color"), "D0D0D0")
            tcPr.append(borders)
    return tbl

def insert_fig(slide, name, left, top, width, caption=None):
    fp = os.path.join(FIGS, name)
    if os.path.exists(fp):
        slide.shapes.add_picture(fp, left, top, width=width)
    if caption:
        tb = slide.shapes.add_textbox(left, top + Inches(3.3), width, Inches(0.35))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = caption; r.font.size = Pt(9); r.font.color.rgb = GRAY75
        r.font.name = "Microsoft YaHei"

def content_slide(slide, title_text, bullets, fig_name=None, sub=None, section="", num=1, bullet_fs=13):
    _set_slide_bg(slide, GRAYF5)
    add_bottom_bar(slide); add_page_no(slide, num); add_section_tag(slide, section)
    add_title(slide, title_text)
    if sub: add_subtitle(slide, sub)
    has_fig = fig_name is not None
    bw_val = 6.8 if has_fig else 12.0
    cy_val = 1.5 if not sub else 1.8
    ch_val = 7.5 - cy_val - 0.7
    add_rect(slide, Inches(0.55), Inches(cy_val), Inches(bw_val), Inches(ch_val), fill_color=WHITE)
    tx = Inches(0.85); ty = Inches(cy_val + 0.2); tw = Inches(bw_val - 0.6); th = Inches(ch_val - 0.4)
    if bullets:
        add_bullets(slide, bullets, tx, ty, tw, th, fs=bullet_fs)
    if has_fig:
        fx = Inches(7.75)
        fy = Inches(cy_val + 0.25)
        fw = Inches(5.0)
        if "fig6_progression" in (fig_name or ""): fw = Inches(5.0)
        insert_fig(slide, fig_name, fx, fy, fw)

def title_slide(slide, title_text, num=1, section=""):
    _set_slide_bg(slide, GRAYF5)
    add_bottom_bar(slide); add_page_no(slide, num); add_section_tag(slide, section)
    add_rect(slide, Inches(0.55), Inches(2.5), Inches(12.2), Inches(2.5), fill_color=WHITE)
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(11.3), Inches(1.9))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title_text; r.font.size = Pt(32); r.font.color.rgb = PRIMARY; r.font.bold = True
    r.font.name = "Microsoft YaHei"

# ═══════ 逐页 ═══════
def build():
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H
    blank = prs.slide_layouts[6]

    # ── P1 封面 ──
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, PRIMARY)
    add_rect(s, Inches(8.5), Inches(0), Inches(5), Inches(7.5), fill_color=RGBColor(0x21, 0x36, 0x5A))
    add_rect(s, Inches(4.2), Inches(2.7), Inches(1.8), Inches(0.01), fill_color=ACCENT)
    tb = s.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(6.8), Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = "基于开源多模态大模型的"; r.font.size = Pt(30); r.font.color.rgb = WHITE; r.font.name = "Microsoft YaHei"
    p = tf.add_paragraph(); r = p.add_run(); r.text = "图像质量评估 Agent 框架"; r.font.size = Pt(36); r.font.color.rgb = WHITE; r.font.bold = True; r.font.name = "Microsoft YaHei"
    p = tf.add_paragraph(); p.space_after = Pt(20)
    r = p.add_run(); r.text = "Unified Agent Framework for Blind IQA"; r.font.size = Pt(15); r.font.color.rgb = ACCENT; r.font.name = "Arial"
    p = tf.add_paragraph(); p.space_after = Pt(28)
    r = p.add_run(); r.text = "张丞  ·  北京邮电大学"; r.font.size = Pt(18); r.font.color.rgb = WHITE; r.font.name = "Microsoft YaHei"
    p = tf.add_paragraph()
    r = p.add_run(); r.text = "指导教师：何老师    验收汇报  ·  2026.07.29"; r.font.size = Pt(13); r.font.color.rgb = GRAY88; r.font.name = "Microsoft YaHei"

    # ═══ 第一章：任务定义与约束（3页）═══
    title_slide(prs.slides.add_slide(blank), "一　任务定义与合规约束", 2, "背景")
    content_slide(prs.slides.add_slide(blank), "1.1  任务目标与Backbone选择", [
        "任务：构建图像质量评估 Agent 框架，不训练/微调 Backbone，仅通过 Skill（Prompt）挖掘大模型的质量评估能力",
        "评价目标：在 KonIQ-10k 验证集与 SPAQ 测试集上取得尽可能高的 SRCC 与尽可能低的 MAE",
        "Backbone：Qwen3-VL-32b-instruct（DashScope API 调用），温度=0 确定性协议，全流程 SHA256 磁盘缓存可复现",
        "对照模型：Qwen3-VL-8b-instruct，用于跨规模对比与门控矩阵迁移实验",
        "评测集规模：KonIQ Val 2,015 张（1–5 整数 MOS）/ SPAQ Test 1,125 张（0–10 连续 MOS）",
        "评测指标：SRCC（斯皮尔曼等级相关系数，主指标）· MAE（平均绝对误差，副指标）· PLCC（皮尔逊线性相关，参考）",
    ], section="背景", num=3)
    content_slide(prs.slides.add_slide(blank), "1.2  合规红线（任务书 §4.1–§4.5）", [
        "§4.1  禁止训练：Backbone 零训练、零微调；任何 Tool / 辅助模型 / 外部模块均不得使用评测数据训练",
        "§4.3  唯一可优化范围：Router/Decision 层的三项职能——评估规则选择 · 多规则冲突裁决 · 评估结果解释生成",
        "§4.4  禁止数据集捷径：数据集名 / Image ID / 文件名 / MOS / 来源分数 / 失真类型 / 失真等级 / 划分信息 → 一律不进在线输入",
        "§4.5  MOS使用边界：仅允许用于离线 Oracle 标签构造 · 离线性能评测 · 离线误差分析 —— MOS 及其衍生信息不进入 Skill/Router/LLM 输入",
        "本文：全程每轮末各读一次 MOS（计算主表），均在预注册文档冻结之后执行；另有一次考后诊断臂读取（依 §4.5“离线性能评测”条款），未用于任何训练或调参",
        "训练集像素（无 MOS）的使用经指导教师邮件书面批准，仅用于自监督信号构造",
    ], section="背景", num=4)

    # ═══ 第二章：方法概述（4页）═══
    title_slide(prs.slides.add_slide(blank), "二　方法概述", 5, "方法")
    content_slide(prs.slides.add_slide(blank), "2.1  系统三层架构", [
        "Backbone 层：Qwen3-VL-32b-instruct，温度=0，全调用 SHA256 摘要哈希缓存——任意实验一键复现、零新增调用",
        "Skill 技能库层：3–5 个评估专家（技术质量 / 整体印象 / 内容完整性，首轮含美学与自然度），每个专家含 7 个可独立开关的提示词组件",
        ("提示词组件：人设 → 维度定义 → 检查清单 → 评估程序 → 文字等级 → 锚点评分带 → 结构化输出契约", GRAY75),
        "Router/Decision 层：全部可学习参数的唯一载体（3 项职能）——7 维像素特征 + 软门控矩阵 W(3×7) 逐图分配话语权",
        ("监督信号全部来自训练集像素自衍生：合成失真阶梯 + 两两比较 Bradley-Terry 锦标赛 —— 零 MOS 接触", GRAY75),
    ], "fig1_pipeline.png", section="方法", num=6)
    content_slide(prs.slides.add_slide(blank), "2.2  Skill 技能库：7 组件独立开关", [
        "① 评估人设：定义专家身份——'你是一位技术质量评估专家' / '普通观者整体第一印象'",
        "② 维度定义：该专家关注的评估维度（噪声/模糊/色彩/构图/内容完整度…）",
        "③ 检查清单：逐项排查的具体失真类型与判断标准",
        "④ 评估程序：5 步强制分析流程——扫描→逐项检查→识别主导问题→判断严重性→映射分数",
        "⑤ 文字等级：5 级文字描述——Bad/Poor/Fair/Good/Excellent",
        "⑥ 锚点评分带（借鉴 Q-Align 机制）：等级文字→分数区间映射（如 Excellent → 4.3-5.0）",
        "⑦ 结构化输出契约：JSON 格式——{\"level\": <int>, \"score\": <float>, \"reason\": \"<= 25 words\"}",
        ("七个组件在消融实验中逐层叠加，精准归因每个组件对 SRCC/MAE 的独立贡献（详见第五章困难与发现）", GRAY75),
    ], section="方法", num=7)
    content_slide(prs.slides.add_slide(blank), "2.3  Router/Decision 层三职能", [
        "职能一：评估规则的选择——7 维手工特征 F（OpenCV 提取）→ 门控矩阵 W(3×7) → softmax 归一化 → 逐图 3 专家话语权向量 g（仅 21 个可训练参数）",
        "职能二：冲突裁决——以 BT 锦标赛排行榜为唯一监督信号，成对铰链损失 −log σ(s_A−s_B) 驱动训练；SPAQ 侧含全局评分与 patch 放大镜的软门控裁决",
        "职能三：解释生成——输出带门控权重与最高权重专家理由的裁决说明，忠实度自检（解释中失真须与专家证据一致）",
        "门控标准：BT 排行榜须在留出对决上超过等权 +0.02；逐图动态融合须超静态权重 +0.005——未过门控即不上场（诚实机制）",
        "额外裁决（SPAQ）：patch 分辨率裁切——1024×1024 裁切 → S-TECH + S-GLOBAL 评分，解决原图压缩导致的细节丢失",
    ], section="方法", num=8)
    content_slide(prs.slides.add_slide(blank), "2.4  纯像素自监督信号（两种形态，零 MOS）", [
        "形态一：合成失真阶梯 —— 对 KonIQ Train 源图施加 6–7 族程序失真（模糊/噪点/JPEG/缩放/色带等），排序天然已知。用于专家体检、Router 训练场与 CKE 门控",
        "体检发现（关键）：模型对变暗近乎不敏感（端点 0.82），对去饱和/过度锐化呈反向偏好（打更高分 0.28–0.43）——三类失真族被除名，不可用作监督",
        "形态二：两两比较锦标赛 — 2,500 张 Train 图像 → 32,389 场对决，双序消除位置偏置，VLM 判胜负 → Bradley-Terry 凸优化 → 每张图一个标量强度分",
        "科学依据：VLM 绝对评分噪声远大于两两比较能力（SPAQ 失真阶梯：绝对值门控全部 FAIL 而比较门控 0.98/0.88）——协议设计先于模型能力",
        "边密度定律：稀疏链图误杀信号；随机长程边 ~7 条/节点方反超旧分排序——经仿真与两轮真数据三次验证",
        "强度分不依赖任何人工标注——是大模型内部质量排序能力的唯一蒸馏物，构成 Router 训练的全部监督信号",
    ], section="方法", num=9)

    # ═══ 第三章：第一轮探索（5页）═══
    title_slide(prs.slides.add_slide(blank), "三　第一轮：探索性实验与核心发现", 10, "第一轮")
    content_slide(prs.slides.add_slide(blank), "3.1  五臂消融设计", [
        "R1-bare：一句裸问'Rate the overall quality of this image on a scale from 1 to 5' → 单次调用直接输出分数（零修饰纯粹先验）",
        "R1-rich：单次调用 + 完整评分细则（五级等级定义 + 检查清单 + 分析流程）→ 代表提示词工程的净收益",
        "R2：五位专家并行评分 + trimmed mean 截尾融合（去掉最高最低取平均）→ 代表多视角分解的收益",
        "R2.5：R2 + 动态分诊 Router —— 画像器判断主导质量问题 → 敏感度矩阵激活排名前 3 专家 → 逆离散度加权融合",
        "R3：R2.5 + CKE 自进化规则库（9 条规则，经双门控筛选后注入各专家提示词尾部）→ 代表自监督迭代的收益",
        ("全部调用温度=0、缓存落盘，五臂同批完成——任何排名差异均可归因至唯一的组件差异", GRAY75),
    ], section="第一轮", num=11)
    content_slide(prs.slides.add_slide(blank), "3.2  首轮主表与分析", [
        "提示词工程在 KonIQ 上产出全场最大单步增益 +0.057 —— 网络野生照片失真多样，系统检查清单补齐模型盲区",
        "同一套详细提示词在 SPAQ 上反而轻退 −0.014 —— 手机 ISP 照片失真集中、模型先验已覆盖，细则引入偏见",
        "等权多专家融合（R2）在两个数据集上均劣于单专家 —— 五位专家同源（同一 VLM），分数相关 0.7–0.9，等权融合实则掺入噪声",
        "R2.5 动态分诊挽回约一半稀释（0.606→0.619），'选谁评比怎么平均更关键' —— 诊断的价值 > 融合的价值",
        "R3 CKE 规则——内部指标全改善（阶梯单调性 +0.019 / B-C 分歧 −0.07），外部 SRCC 纹丝不动（+0.002/–0.001）",
    ], section="第一轮", num=12)
    content_slide(prs.slides.add_slide(blank), "3.2  首轮主表", [], section="第一轮", num=13)
    # Manual table placement for this slide
    s13 = prs.slides[-1]
    hdr = ["臂", "机制", "KonIQ SRCC", "KonIQ MAE", "SPAQ SRCC", "SPAQ MAE"]
    rows = [hdr,
            ["R1-bare", "裸问直接评分", "0.577", "0.463", "0.881", "0.893"],
            ["R1-rich", "单专家+完整细则", "0.633", "0.688", "0.867", "1.042"],
            ["R2", "五专家+截尾融合", "0.606", "0.677", "0.859", "0.971"],
            ["R2.5", "+动态分诊Router", "0.619", "0.533", "0.860", "1.014"],
            ["R3", "+CKE自进化规则库", "0.621", "0.499", "0.861", "1.008"]]
    make_table(s13, rows, Inches(0.55), Inches(1.55),
               [Inches(1.2), Inches(3.2), Inches(1.35), Inches(1.3), Inches(1.35), Inches(1.3)])
    add_bullets(s13, [
        ("KonIQ 增益链：R1b 0.577 → R1r 0.633（+0.057，prompt工程）→ R2.5 0.619（Router优化）→ R3 0.621（CKE零收益）", GRAY75),
        ("SPAQ 天花板现象：R1b 0.881 已接近——复杂机制无处可增 → 原因在第五轮SPAQ协议瓶颈分析中揭示", GRAY75),
    ], Inches(0.85), Inches(4.2), Inches(11.5), Inches(2.5))

    content_slide(prs.slides.add_slide(blank), "3.3  CKE 自进化规则库：内部全改善、外部零收益（锚相关性定律雏形）", [
        "CKE 流程：B 路线（分析型打分）与 C 路线（两两对比+BT 拟合）在 1,000 张工作集上求分歧 → 裁判 LLM 提炼规则 → 双门控筛选 → 四轮收敛 9 条",
        "内部指标：阶梯单调性 0.455→0.474（+0.019）、B-C 分歧 0.402→0.330（−0.07）——全部改善",
        "外部指标：KonIQ SRCC R2.5 0.619 → R3 0.621（+0.002）——纹丝不动。最昂贵的组件（¥69）产出近零收益",
        "根因分解：① CKE 的驱动信号（自洽性 + 合成失真阶梯）优化的是模型内部一致性，而非与人群分的对齐；② 合成失真仅覆盖 4 族程序失真，与真实野生失真上的 MOS 相关性弱",
        ("这一阴性结果是全项目最重要的发现——催生了核心假设：零训练自监督 IQA 的上限由监督信号与真实目标的相关性决定（锚相关性定律）", ACCENT),
    ], section="第一轮", num=14)

    # ═══ 第四章：第二轮预注册验证（4页）═══
    title_slide(prs.slides.add_slide(blank), "四　第二轮：预注册验证与机制深化", 15, "第二轮")
    content_slide(prs.slides.add_slide(blank), "4.1  预注册方法论的引入", [
        "第一轮的定性发现需要转化为可检验的定量假设。第二轮引入预注册流程——先冻结先验假设（预注册-R4R5 文档），再执行实验、读取 MOS 检验",
        "预注册假设 H1：同一融合机制在正确的（同分布 BT）信号上，产出远超旧梯子信号的收益",
        "预注册假设 H2：SPAQ 三槽位软门控协议路由 + 释义投票，可提升 SRCC 并降低 MAE",
        "预注册假设 H3：R4-KonIQ MAE ≤ 第一轮静态融合最优 MAE（0.533）",
        ("门控通过标准：BT 排行榜留出一致率 > 等权 +0.02；静态门控权重大于等权 +0.01；期望类评分排序不降且平滑度提升", GRAY75),
        ("任何假设不成立 → 主表如实呈现，阴性结果与阳性同等写入验收材料（延续第一轮的证伪传统）", GRAY75),
    ], section="第二轮", num=16)
    content_slide(prs.slides.add_slide(blank), "4.2  R4：同分布 BT 监督取代旧梯子信号", [
        "R4 与 R2 使用同一份缓存的 5 专家评分，唯一改变：融合器由 trimmed mean 换为以同分布 BT 锦标赛为监督的静态学习权重",
        "方法：对 KonIQ Train 2,500 张组织 BT 锦标赛 → 每张图一个标量强度分 → 以强度分差为目标训练融合权重 → 留出验证一致率",
        "结果（预注册 H1 命中）：KonIQ SRCC 从 R2 的 0.606 跃升至 0.668（+0.062）——而旧梯子信号上同一学习权重方法仅产出 +0.007",
        "归因：同一融合机制、同一代码、不同信号——'信号更换 → 收益量级跃迁'，将锚相关性定律从定性发现升级为定量验证",
        "交叉验证：4 族失真留一交叉验证，拟合权重 vs 等权差值仅 0.005–0.009，验证无过拟合",
        ("KonIQ 侧未设 R5：图像尺寸恒为 512×384，无 patch 放大需求；R4 的 BT 监督路由已在同域取得强证据。", GRAY75),
    ], section="第二轮", num=17)
    content_slide(prs.slides.add_slide(blank), "4.3  R5：SPAQ 三槽位软门控与释义投票", [
        "SPAQ 架构（R4/R5 共用）：三个评分槽位独立输出分数，由软门控按图像特征逐图分配话语权重",
        "槽位一（bare）：裸问直接评分（1 条 prompt），权重约 0.5 —— 零点基准",
        "槽位二（rich-para3）：细则释义评分（3 条 prompt），权重约 0.34 —— 细则收益",
        "槽位三（multi-patch）：原分辨率裁切（1024×1024）多专家评分，权重约 0.16 —— 细节放大",
        "新增组件——裸问释义投票：4 条同义问法取均值（smooth quantization），单位数增多、评分更平滑",
        "预注册 H2 半中：SPAQ SRCC 0.891（+0.006 命中），MAE 0.944（劣于 R4 的 0.908，H3 未中）→ 零点偏置部分抵消路由收益",
    ], section="第二轮", num=18)

    # ═══ 第五章：第三轮统一框架（6页）═══
    title_slide(prs.slides.add_slide(blank), "五　第三轮：统一框架收敛", 19, "第三轮")
    content_slide(prs.slides.add_slide(blank), "5.1  设计收敛：从五臂到统一骨架", [
        "在锚相关性定律指导下，第三轮（R6）完成三项关键收敛：",
        "收敛一：专家瘦身——首轮证伪的美学（S-AESTH）与自然度（S-NATURAL）被移除，保留技术/整体/内容三个专家",
        "收敛二：统一骨架——两个数据集、两个模型规模共用同一套代码，仅参数不同（α 与门控是否激活）",
        "收敛三：三门控柱收敛为单一路由——评估规则选择 → 冲突裁决 → 解释生成，全部由 W(3×7) 门控矩阵承载",
        "KonIQ 配置：逐图动态软门控（BT 监督训练）→ 融合分 + 裸问释义投票，α = 0.6（300 张 Train pilot 按 τ 准则扫描）",
        "SPAQ 配置：三槽位软门控（门控未过一致性门槛 → 自动回退等权）→ 融合分 + 裸问释义投票，α = 0.3（200 张扫描）",
    ], "fig1_pipeline.png", section="第三轮", num=20)
    content_slide(prs.slides.add_slide(blank), "5.2  门控矩阵训练：21 个参数的全链路", [
        "初始化：W = 全零矩阵（3 行 × 7 列），7 维标准化特征包含 lap_var / noise / colorful / bright / logpix / aspect / spread",
        "每步：随机抽 128 对图 → 正向计算 W×F → softmax 得话语权 g → 融合分 s_fus = g·s → 成对铰链损失 −log σ(s_A−s_B) → 解析梯度 ∂L/∂W → 更新",
        "训练曲线：800 步收敛，损失 0.46→0.40，参数范数 ‖W‖ 0.0→3.8，CPU 数秒完成——全程不调大模型、不查 BT 排行榜（只在监督构造阶段使用）",
        "留出检验：2,000 对图上动态融合一致率 80.9%，等权融合 78.0%，动态领先 2.9 个百分点——门控学到了可泛化的路由规律",
        "收敛后门控矩阵冻结——推理阶段：每张图仅做一次矩阵乘法（21 个乘法 + softmax + 3 个乘法加权），毫秒级，不调任何大模型",
    ], "fig2_loss.png", section="第三轮", num=21)
    content_slide(prs.slides.add_slide(blank), "5.3  门控矩阵的可解释性与真实算例", [
        "训练后 W 矩阵高度可读——分歧度（spread）列为绝对主导特征：S-TECH +0.89、S-CONTENT −0.84",
        "→ 门控的决策逻辑：'专家意见不一致时，信任技术专家，抑制内容专家的权重'",
        "logpix / aspect 权重自动归零——训练集尺寸恒为 512×384，无判别信息，训练算法自主忽略（无需人工特征选择）",
        "真实推理算例——某图分歧度高出 +2.47σ：logit_TECH +2.986 / logit_GLOBAL −0.215 / logit_CONTENT −2.771",
        "→ softmax 归一化后：TECH 95.8% / GLOBAL 3.9% / CONTENT 0.3% —— 技术专家几乎独揽话语权",
        "→ 融合分 1.90（等权融合会被 S-GLOBAL 拉高至 3.53，接近翻倍失真）—— 门控有效抑制分布外专家干扰",
    ], "fig3_w.png", section="第三轮", num=22)
    content_slide(prs.slides.add_slide(blank), "5.4  预注册假设检验（冻结→实验→检验→归因）", [
        "H1'：R6-KonIQ SRCC > R4（0.668）——命中：0.734（+0.065）。动态融合 +0.026 一致率与 α 混合 +0.045 τ 全额兑现",
        "H2'：R6-SPAQ SRCC ≥ R5（0.891）且 MAE < R5（0.944）——半中：SRCC 0.893 命中，MAE 0.945 差 0.001 未中——SPAQ 零点偏置为结构性边界",
        "H3'：R6-KonIQ MAE ≤ R4（0.554）——命中：0.491（−0.063）。裸问释义投票的零点校准作用被实证——bare 通道的偏置远小于融合通道",
        "SPAQ 门控未过一致性门槛（动态融合一致率低于等权融合 +0.005）→ 框架自动回退等权融合——内置的诚实机制：门控不达标、绝不上场",
        ("三条假设的中与不中全部写入主表并附归因解释——预注册框架的核心价值：保证可证伪性，阴性结果的边界条件与阳性结果同等清晰", ACCENT),
    ], section="第三轮", num=23)
    content_slide(prs.slides.add_slide(blank), "5.5  考后诊断臂（post-hoc, 未预注册, 不参与主线结论）", [
        "R1-anchor（F-024）：仅保留裸问+等级锚点评分带，剥除专家人设/检查清单/JSON 程序——用于分解 R1-bare→R1-rich 增益的来源",
        "→ KonIQ 裸分 0.577 → 仅锚点 0.641 → 完整 rich 0.633 —— 排序增益几乎全部来自锚点层，专家程序层贡献为微负 ≈ 与同源稀释互为印证",
        "R6-offanchor（F-025）：R6 结构完全不变，仅砍 _SCALE_BLOCK 锚点评分带——用于测试锚点在多专家融合中的角色",
        "→ KonIQ 0.729（略劣于 R6 的 0.734）、SPAQ MAE 从 0.945 崩塌至 1.208 —— S-TECH 均值低 1.45 分，跨专家分数不在同一天平",
        ("综合结论：锚点单独有害（裸分场景劣化 MAE），锚点联合有益（多专家融合中是跨维度校准的公共标尺）——同一组件在不同架构深度扮演相反角色", ACCENT),
    ], section="第三轮", num=24)

    # ═══ 第六章：综合结果与归因（4页）═══
    title_slide(prs.slides.add_slide(blank), "六　综合结果与归因分析", 25, "结果")
    # P26 -- main table
    s26 = prs.slides.add_slide(blank)
    _set_slide_bg(s26, GRAYF5); add_bottom_bar(s26); add_page_no(s26, 26); add_section_tag(s26, "结果")
    add_title(s26, "6.1  三轮综合主表")
    hdr = ["轮次", "臂", "KonIQ SRCC", "KonIQ MAE", "SPAQ SRCC", "SPAQ MAE"]
    rows = [hdr,
            ["一", "R1-bare（裸问基线）", "0.577", "0.463", "0.881", "0.893"],
            ["一", "R1-rich（+完整细则）", "0.633", "0.688", "0.867", "1.042"],
            ["一", "R2（5专家+截尾融合）", "0.606", "0.677", "0.859", "0.971"],
            ["一", "R2.5（+动态分诊）", "0.619", "0.533", "0.860", "1.014"],
            ["一", "R3（+CKE规则9条）", "0.621", "0.499", "0.861", "1.008"],
            ["二", "R4（BT监督路由）", "0.668", "0.554", "0.885", "0.908"],
            ["二", "R5（SPAQ门控+投票）", "—", "—", "0.891", "0.944"],
            ["三", "R6（统一框架）", "0.734", "0.491", "0.893", "0.945"],
            ["*", "R1-anchor（考后诊断）", "0.641", "0.544", "0.894", "0.807"],
            ["*", "R6-offanchor（考后诊断）", "0.729", "0.576", "0.894", "1.208"]]
    make_table(s26, rows, Inches(0.55), Inches(1.55),
               [Inches(0.6), Inches(3.0), Inches(1.35), Inches(1.3), Inches(1.35), Inches(1.3)])
    # KPI boxes for R6
    add_kpi_box(s26, Inches(9.6), Inches(1.5), Inches(3.1), "KonIQ SRCC (R6 vs R1b)", "0.734", "+0.157", BLUE_A)
    add_kpi_box(s26, Inches(9.6), Inches(2.55), Inches(3.1), "KonIQ MAE (R6)", "0.491", "vs R1b 0.463", ACCENT)
    add_kpi_box(s26, Inches(9.6), Inches(3.6), Inches(3.1), "SPAQ SRCC (R6)", "0.893", "+0.012 vs R1b", BLUE_A)
    add_kpi_box(s26, Inches(9.6), Inches(4.65), Inches(3.1), "越过 8B 零样本锚点", "0.734 > 0.729", "Tool-IQA ref.", GREEN_A)

    content_slide(prs.slides.add_slide(blank), "6.2  归因分析：每一分增益对应具体决策", [
        "KonIQ 总增益 +0.157 的逐步分解（R1b 0.577 → R6 0.734）：",
        "段一：R1b→R1r +0.057 —— 提示词工程增益，几乎全部来自锚点层（仅锚点即 0.641），专家程序层无独立排序贡献（考后 R1-anchor 探针证实）",
        "段二：R2→R4 +0.062 —— 监督信号更换：旧梯子信号上同一融合机制产出 +0.007，换为同分布 BT 锦标赛后同一代码产出 +0.062（锚相关性定律直接验证）",
        "段三：R4→R6 +0.065 —— 融合从静态权重演进为逐图动态门控（+0.026 一致率）+ 裸问释义投票独立零点校准（+0.045 τ）",
        ("SPAQ 总增益 +0.012（R1b 0.881→R6 0.893）：软门控路由+0.004 · patch 放大镜与释义投票+0.006 · 裸问投票替换+0.002 —— 天花板高、余量小", GRAY75),
        "交叉验证：在统一框架骨架内同步跑了 8B 移栽实验，门控矩阵跨模型迁移零新增 API 调用，SRCC 0.759→0.776（详见 6.4 节）",
    ], section="结果", num=27)
    content_slide(prs.slides.add_slide(blank), "6.3  与外部工作对比（评测协议不同，仅供参考）", [
        "Tool-IQA（arXiv:2606.16082）在同一 Backbone 族（Qwen3-VL-8B）上报告了三个配置：",
        "零样本直接评分：KonIQ 0.729 / SPAQ 0.856 —— 本文 32B 在 KonIQ 上以 0.734 越过此锚点",
        "工具增强（training-free）：KonIQ 0.732 / SPAQ 0.866 —— 本文 SPAQ 0.893 大幅超越（+0.027）",
        "GRPO 训练后（本任务 §4.1 不可用）：KonIQ 0.825 / SPAQ 0.898 —— 本文 SPAQ 0.893 逼近训练后水平（−0.005）",
        "需注意：评测协议不同（他们用 8B、不同 prompt、不同划分），直接数值对比仅供参考",
        ("本文核心优势在于：零训练、零 MOS、仅 21 参数可训练、全部结果缓存可复现——方法论可证伪性远高于数值本身", ACCENT),
    ], section="结果", num=28)
    content_slide(prs.slides.add_slide(blank), "6.4  门控矩阵跨模型迁移（零新增 API 调用）", [
        "实验：将 32B Backbone 上训练的 W(3×7) 门控矩阵，原样作用于 8B Backbone 的同一三专家池——纯矩阵乘法、零 API 调用（缓存命中 12,084 次）",
        "8B 纯裸问探针：SRCC 0.679 —— 模型规模对人群趋势把握的基准能力",
        "8B + 等权门控：SRCC 0.759 —— 三专家集成本身提供 +0.080 增益（相对裸问）",
        "8B + 迁移 32B 门控：SRCC 0.776 —— 额外 +0.017 且反超 32B 自身（0.734）",
        "结论：门控学到的是'图像属性 → 专家可信度'的通用映射，与 Backbone 规模解耦——锚相关性定律的跨模型稳定性由此实证",
        "反超现象的解释：8B 专家分与 32B 专家分的排序结构高度相似，但 8B 等权融合时被个别异常专家拉偏得更厉害 → 门控抑制作用更显著",
    ], "fig7_transfer.png", section="结果", num=29)

    # ═══ 第七章：困难失败与关键发现（5页）═══
    title_slide(prs.slides.add_slide(blank), "七　困难、失败与关键发现", 30, "发现")
    content_slide(prs.slides.add_slide(blank), "7.1  CKE 自进化规则库（无效组件，¥69 / +0.001 SRCC）", [
        "设计初衷：B/C 两范式分歧 → 裁判 LLM 提炼规则 → 双门控筛选 → 注入专家提示词 → 期望提升外部 SRCC",
        "执行：四轮迭代，+3/+3/+3/+0 收敛于 9 条规则，内部指标全部改善（阶梯单调性 +0.019 / 分歧 −0.07），裁判 LLM 一轮 ¥17",
        "外部结果：KonIQ SRCC +0.001、SPAQ −0.001 —— 近乎为零。最昂贵也最无收益的组件",
        "根因：两个驱动信号（B-C 自洽性 + 合成失真阶梯单调性）优化的是模型内部一致性，而非与人群分的对齐——自洽≠对齐",
        "对照证据：3DrawAgent 能靠自监督 CKE 获益，因 CLIP 作为外部锚与语义一致性天然相关；我们的阶梯锚仅覆盖 4 族合成失真——相关性弱",
        ("该失败是锚相关性定律的第一块基石：信号与目标的相关性决定上限。CKE 的全部代码与规则库封存备查", ACCENT),
    ], section="发现", num=31)
    content_slide(prs.slides.add_slide(blank), "7.2  同源多专家稀释（无效设计）", [
        "五位专家同源（同一 VLM），分数相关性 0.7–0.9 → 等权融合使 KonIQ SRCC 从 0.633 降至 0.606（−0.027 稀释）",
        "在旧梯子信号上拟合的学习权重接近等权（+0.007，4 族交叉验证无过拟合）——同一 VLM 的多视角分解无法创造新的判别信息",
        "'解释的多样性 ≠ 预测的多样性' —— 不同提示词诱导出了不同的语言表达，但底层视觉理解共享同一套参数",
        ("该证伪催生了第三轮的专家瘦身（5→3）、确立了'用 BT 相对比较信号而非同源融合重加权'的技术方向", GRAY75),
        "后续验证：第三轮的三专家池设计受益于此——美学与自然度专家在交叉验证中始终权重最低，移除后无任何 SRCC 损失",
    ], section="发现", num=32)
    content_slide(prs.slides.add_slide(blank), "7.3  分布外监督失效与失真族反向偏好（缺陷组件）", [
        "合成失真阶梯仅覆盖 4 族程序失真（模糊/噪点/JPEG/变暗），与真实野生失真分布不一致——分布外锚天然弱相关",
        "体检发现 3 族无效失真：变暗不敏感（端点准确率 0.82，相邻档仅 0.41）；去饱和反向偏好（失真版打更高分 0.28）；过度锐化反向偏好（0.43）",
        "BT 锦标赛的边密度定律：稀疏链图误杀信号——仿真显示随机长程边约 7 条/节点时方反超旧分排序",
        "SPAQ 绝对评分塌方：6 族失真门控全部 FAIL（模糊仅 0.65），而同一批图的两两比较却达 0.98/0.88——相差一个量级",
        "这一发现直接确立了第三轮'排序优先、比较先行'的协议路线——协议设计先于模型能力",
    ], section="发现", num=33)
    content_slide(prs.slides.add_slide(blank), "7.4  MAE 零点偏置（结构性边界，在 §4.5 约束下原则性不可达）", [
        "三轮主线各臂 MAE 均未超过裸分 0.463/0.893——偏置是 MAE 的主项，而非方差",
        "病理分解：模型先验的内容依赖零点摆幅约 19pp，标注人群约 6pp——对野生网图宽容、对手机照片苛刻，各超调一半",
        "提示词组件归因：KonIQ 裸分偏移 +0.220 → 加锚点层 +0.453 → 加专家程序层 +0.581——约 65% 附加偏置来自锚点层",
        "损益变号（关键）：SPAQ 上同一锚点表将零点拉近 MOS，MAE 0.807 反超裸分 0.893 —— 盲写先验不可能同时命中两数据集",
        "任何常数平移 = MOS 校准——§4.5 禁止。该偏置在零 MOS 设定下原则性不可达，但内容条件化零点为开放残留课题",
        ("KonIQ PLCC 最高（0.778）：'内容赢了、刻度输了'——排序与校准的能力分离，框架内部自洽", ACCENT),
    ], "fig5_ablation.png", section="发现", num=34)
    content_slide(prs.slides.add_slide(blank), "7.5  锚点的双重角色（单问 vs 融合的镜像证据）", [
        "R1-anchor（单问场景）：锚点单独有害——裸分 MAE 0.463 → 加锚点 0.544 → 加完整专家程序 0.688（±30pp 恶化）",
        "R6-offanchor（融合场景）：锚点联合有益——砍锚后 KonIQ 0.734→0.729（+6pp 可控）SPAQ MAE 0.945→1.208（崩盘 +263pp）",
        "解剖 SPAQ：offanchor S-TECH 均值比锚定版低 1.45 分，模型对手机照片裸技术先验（~3.2/10）比整体先验（~5.2/10）严苛 2 级",
        "→ 无锚点表时跨专家分数不在同一天平上，融合器无法区分图间差异——锚点是专家间的公共标尺，而非独立存在的评分知识",
        ("综合：同一组件在架构浅层是纯成本（框死零点），在架构深层是必需品（跨专家公共标尺）——R6 恰好处于两条曲线的交点", ACCENT),
        ("这是本项目对'零训练下如何设计评估架构'给出的最精确量化回答", ACCENT),
    ], "fig8_dist.png", section="发现", num=35)

    # ═══ 第八章：讨论与结论（3页）═══
    title_slide(prs.slides.add_slide(blank), "八　讨论与结论", 36, "结论")
    content_slide(prs.slides.add_slide(blank), "8.1  锚相关性定律的完整闭合", [
        "第一轮（发现）：CKE 零收益证伪'自洽性驱动自监督改进'→ 同源多专家 +0.007 近零增益证伪'旧信号上重加权可行'→ 确立'信号相关性决定上限'",
        "第二轮（验证）：同一融合机制在旧梯子信号上 +0.007 vs 同分布 BT 信号上 +0.062——同一代码、不同信号、收益差近一个量级——定性命题升级为定量定律",
        "第三轮（放大）：扩容锦标赛至 32,000 场对决 → 引入逐图动态门控（21 参数）→ 增设独立零点校准通道（bare 投票）→ KonIQ SRCC 0.577→0.734 / SPAQ 0.881→0.893",
        "跨模型迁移：32B 门控 → 8B 专家，SRCC 0.776 反超原主——信号质量高于模型规模",
        ("三轮 = 发现定律 → 验证定律 → 放大定律，构成完整的'假设驱动实验→预注册检验→门控放大→跨模型泛化'闭环", ACCENT),
    ], section="结论", num=37)
    content_slide(prs.slides.add_slide(blank), "8.2  本文贡献", [
        "方法贡献：在零训练、零 MOS、仅 Router 层可优化的严苛约束下，构建了两数据集 × 两模型规模复用的统一 Agent 框架（温度=0、全缓存可复现）",
        "理论贡献：提出并验证了锚相关性定律——零训练自监督 IQA 的上限由监督信号与真实目标的相关性决定——提供从定性发现到定量验证到放大实证的完整证据链",
        "方法论贡献：预注册加固研究的可归因性与可证伪性——CKE 零收益、同源多专家稀释、分布外锚失效、MAE 零点偏置等全部阴性结果如实记录并给出根因解释",
        "工程贡献：设计并稳定运行了一套 10 万级 API 调用的可复现实验系统，克服限流/进程残留/缓存缺陷/断电/批量事故等六项工程困难",
    ], section="结论", num=38)
    content_slide(prs.slides.add_slide(blank), "8.3  局限与展望", [
        "局限一：评价上限受 Backbone 能力约束——32B vs 8B 差距（0.734 vs 0.776 迁移后）即模型能力阈的实证",
        "局限二：SPAQ 门控未过一致性门槛——通用像素特征对该域判别力不足，门控在此域价值接近于维持而非提升",
        "局限三：MAE 零点偏置在 §4.5 约束下原则性不可达——内容条件化零点为开放的残留课题",
        "局限四：BT 锦标赛成本随图像量平方增长——大规模应用需采样策略（当前 ~7 边/节点的密度定律可指导成本控制）",
        "横向推广：框架不依赖特定任务结构——可迁移至视频质量评价、生成图像美学评估、医学影像质控等以人群主观分为标准的外部评价场景",
        "纵向深入：门控矩阵在线增量更新 · 锦标赛分层采样策略 · 分块级路由定位局部退化 · '裸问校零'从经验设计提升为可证明的校准机制",
    ], section="结论", num=39)

    # ── P40 结尾页 ──
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, PRIMARY)
    add_rect(s, Inches(0), Inches(0), Inches(0.06), H, fill_color=ACCENT)
    add_rect(s, Inches(0), H - Inches(0.06), W, Inches(0.06), fill_color=ACCENT)
    add_page_no(s, 40)
    tb = s.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = "谢谢各位老师，请批评指正"; r.font.size = Pt(36); r.font.color.rgb = WHITE; r.font.bold = True; r.font.name = "Microsoft YaHei"
    p = tf.add_paragraph(); p.space_after = Pt(20)
    r = p.add_run(); r.text = "张丞  ·  北京邮电大学  ·  2026.07.29"; r.font.size = Pt(18); r.font.color.rgb = ACCENT; r.font.name = "Microsoft YaHei"
    p = tf.add_paragraph(); r = p.add_run(); r.text = "材料与代码：github.com/...  ·  docs/findings.md（F-001~F-025 全部证据链）"; r.font.size = Pt(12); r.font.color.rgb = GRAY88; r.font.name = "Microsoft YaHei"

    # ── 备用页 2 张 ──
    title_slide(prs.slides.add_slide(blank), "（备用一）工程困难实录", 41, "备用")
    content_slide(prs.slides.add_slide(blank), "工程六难", [
        "① DashScope 限流 ~600 RPM：并发 24→16、指数退避 8 次重试、单点容错 + 磁盘缓存断点续跑——保障 10 万级调用的零丢失",
        "② 僵尸运行：Windows TaskStop 只杀父进程——子进程跑完全部 10 臂，账单单差由此对平（¥160 vs ¥97）",
        "③ 缓存键缺失 temperature：'重跑'实为静默重放——生成参数并入缓存键后当场检出，既有 2 万条缓存零污染",
        "④ 批量操作三事故：字典误传协程致万场对决瞬时全灭 · 新旧节点编号撞车 · id 平移不幂等污染 10,694 条边——手术清除+全员标记+立规'批量先小样、写入必幂等'",
        "⑤ 两次断电：增量落盘补丁 + 缓存断点续跑——三次续跑零数据丢失",
        "⑥ 内容审查：2 张图触发 API 安全过滤（KonIQ / SPAQ 各 1 张）——排除并文档化",
    ], section="备用", num=42, bullet_fs=12)

    content_slide(prs.slides.add_slide(blank), "（备用二）合规自检清单", [
        "在线输入白名单核对：图像像素 + 七个可独立开关的提示词组件 + 温度=0 确定性协议——仅此三项，不含任何禁止信息",
        "数据集名：在 Router/LLM/Skill 的输入中零出现——全部由配置文件的 scale key（'koniq'/'spaq'）映射到标尺，无数据集语义",
        "MOS 唯一读取入口：iqa_agent/data.py → load_mos() —— 全程仅在 50_eval.py 中调用，每轮末各读一次（共三次），均在预注册文档冻结之后",
        "训练信号自证：失真阶梯源图来自 KonIQ Train 划分（经邮件批准），BT 锦标赛强度分仅依赖 VLM 两两比较——两份信号均未接触 MOS",
        "缓存完整性：SHA256 摘要 key = (model, prompt_hash, image_hash, temperature) —— 温度参数已嵌入键，避免静默重放",
        "R4/R5 中间产物的有效成分（BT 监督路由、裸问释义投票、patch 放大镜）已全部并入 R6 统一框架——无冗余分支遗留",
    ], section="备用", num=43)

    path = os.path.join(ROOT, "docs", "汇报_detail.pptx")
    prs.save(path)
    print("->", path)

if __name__ == "__main__":
    build()
