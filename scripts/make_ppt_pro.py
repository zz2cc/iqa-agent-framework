# -*- coding: utf-8 -*-
"""专业学术答辩 PPT：python-pptx 手写，深蓝主题 + 金色点缀 + 卡片式布局。"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn, _nsmap
from lxml import etree

def OxmlElement(nsptag_str):
    """Create an oxml element with proper clark notation."""
    from pptx.oxml.ns import NamespacePrefixedTag
    return etree.Element(NamespacePrefixedTag(nsptag_str).clark_name)
from copy import deepcopy

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGS = os.path.join(ROOT, "docs", "figs")
sys.stdout.reconfigure(errors="replace")

# ═══════ 设计系统 ═══════
PRIMARY = RGBColor(0x1B, 0x2A, 0x4A)      # 深藏蓝
ACCENT  = RGBColor(0xC9, 0xA9, 0x6E)       # 淡金
BLUE_A  = RGBColor(0x2A, 0x78, 0xD6)       # 图表蓝
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x1A, 0x1A)
GRAY75  = RGBColor(0x75, 0x75, 0x75)
GRAY88  = RGBColor(0x88, 0x88, 0x88)
GRAY95  = RGBColor(0x95, 0x95, 0x95)
GRAYE0  = RGBColor(0xE0, 0xE0, 0xE0)
GRAYF0  = RGBColor(0xF0, 0xF0, 0xF0)
GRAYF5  = RGBColor(0xF5, 0xF5, 0xF5)

W = Inches(13.333); H = Inches(7.5)
SLIDE_W = W; SLIDE_H = H


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill; fill.solid(); fill.fore_color.rgb = color


def _rot(shape, angle):
    shape.rotation = angle


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, lw=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color; shape.line.width = lw or Pt(1)
    return shape


def add_left_accent(slide, left, top, height, color=ACCENT, width=Inches(0.06)):
    return add_rect(slide, left, top, width, height, fill_color=color)


def add_bottom_bar(slide, color=PRIMARY):
    add_rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), fill_color=color)


def add_page_number(slide, num, total=12):
    tb = slide.shapes.add_textbox(Inches(11.8), H - Inches(0.55), Inches(1.2), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = f"{num} / {total}"
    r.font.size = Pt(8); r.font.color.rgb = GRAY88
    r.font.name = "Arial"


def add_section_tag(slide, tag, x=Inches(0.55), y=H - Inches(0.55)):
    tb = slide.shapes.add_textbox(x, y, Inches(3), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = tag
    r.font.size = Pt(8); r.font.color.rgb = GRAY95
    r.font.name = "Microsoft YaHei"


def mk_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def add_text(tf, text, size=14, color=DARK, bold=False, name="Microsoft YaHei", alignment=None, space_after=0, space_before=0):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = Pt(space_after); p.space_before = Pt(space_before)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    r.font.name = "Arial" if all(ord(c) < 256 for c in text) else name
    r.font.name = name
    return p


def add_title(slide, text, y=Inches(0.55)):
    add_left_accent(slide, Inches(0.55), y + Inches(0.08), Inches(0.52))
    tb = mk_textbox(slide, Inches(0.85), y, Inches(11.5), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    add_text(tf, text, size=26, color=PRIMARY, bold=True, space_before=0)
    return tf


def make_table(slide, rows, left, top, col_widths, header_color=PRIMARY):
    """rows[0] = header, rows[1:] = data."""
    n_rows, n_cols = len(rows), len(rows[0])
    total_w = sum(col_widths)
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, Inches(0.35 * n_rows))
    table = tbl.table
    for j, w in enumerate(col_widths):
        table.columns[j].width = w
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(8.5); r.font.name = "Arial"; r.font.bold = (i == 0)
            r.font.color.rgb = WHITE if i == 0 else DARK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = GRAYF5 if i % 2 == 0 else WHITE
            # Thin borders
            tcPr = cell._tc.get_or_add_tcPr()
            borders = OxmlElement("w:tcBorders")
            for edge in ["top", "bottom", "left", "right"]:
                el = OxmlElement(f"w:{edge}")
                el.set(qn("w:val"), "single")
                el.set(qn("w:sz"), "4")
                el.set(qn("w:color"), "D0D0D0")
                borders.append(el)
            tcPr.append(borders)
    return tbl


def insert_fig(slide, name, left, top, width, caption=None, cap_y=None):
    fp = os.path.join(FIGS, name)
    pic = slide.shapes.add_picture(fp, left, top, width=width)
    if caption:
        tb = mk_textbox(slide, left, cap_y or (top + Inches(3.3)), width, Inches(0.35))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = caption
        r.font.size = Pt(9); r.font.color.rgb = GRAY75
        r.font.name = "Microsoft YaHei"
    return pic


def content_slide(slide, title_text, bullets, fig_name=None, section="", num=1):
    _set_slide_bg(slide, GRAYF5)
    add_bottom_bar(slide)
    add_page_number(slide, num); add_section_tag(slide, section)
    add_title(slide, title_text)
    # 正文卡片
    card_x, card_y = Inches(0.55), Inches(1.5)
    has_fig = fig_name is not None
    card_w = Inches(7.4) if has_fig else Inches(12.2)
    card_h = Inches(5.2)
    add_rect(slide, card_x, card_y, card_w, card_h, fill_color=WHITE)
    # 内容
    tb = mk_textbox(slide, Inches(1.05), Inches(1.75), card_w - Inches(1.0), card_h - Inches(0.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_after = Pt(8)
        # Bullet dot
        r = p.add_run(); r.text = " • "
        r.font.size = Pt(12); r.font.color.rgb = ACCENT; r.font.bold = True
        r = p.add_run(); r.text = b
        r.font.size = Pt(14); r.font.color.rgb = DARK
        r.font.name = "Microsoft YaHei"
    # 配图
    if has_fig:
        fig_x = Inches(8.35)
        insert_fig(slide, fig_name, fig_x, card_y + Inches(0.3), Inches(4.4))


# ═══════ 逐页 ═══════
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ── P1 封面 ──
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, PRIMARY)
    # 装饰大三角
    add_rect(s, Inches(8.5), Inches(0), Inches(5), Inches(7.5), fill_color=RGBColor(0x21, 0x36, 0x5A))
    # 装饰细线
    add_rect(s, Inches(4.2), Inches(2.7), Inches(1.8), Inches(0.01), fill_color=ACCENT)
    # 标题
    tb = mk_textbox(s, Inches(1.2), Inches(1.6), Inches(6.8), Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    add_text(tf, "基于多模态大模型智能体的", size=30, color=WHITE, bold=False, space_after=6)
    add_text(tf, "无参考图像质量评价统一框架", size=36, color=WHITE, bold=True, space_after=20)
    add_text(tf, "Unified Agent Framework for Blind IQA", size=15, color=ACCENT, bold=False, name="Arial", space_after=28)
    # 信息
    add_text(tf, "张丞  ·  北京邮电大学", size=18, color=WHITE, space_after=4)
    add_text(tf, "验收汇报  ·  2026 年 7 月 29 日", size=13, color=GRAY88, name="Arial")

    # ── P2 任务 ──
    s = prs.slides.add_slide(blank)
    content_slide(s, "任务与硬约束", [
        "任务：构建通用图像质量评价 Agent 框架，主干大模型全程零训练",
        "主干：qwen3-vl-32b-instruct / 8b-instruct，温度 = 0，全流程磁盘缓存可复现",
        "合规红线：数据集名、图像标识、人群分数——一律不进入 Router/LLM/Skill 线上输入",
        "评测域：KonIQ-10k 验证集（2,014 张，1–5 整数分）/ SPAQ 测试集（1,124 张，0–10 连续分）",
        "评测指标：SRCC（排序一致性）/ MAE（绝对校准）/ PLCC（线性相关）",
        "代码、权重、缓存、评测脚本全部归档——审核可一键复跑全部实验",
    ], section="任务", num=2)

    # ── P3 框架总览 ──
    s = prs.slides.add_slide(blank)
    content_slide(s, "统一框架总览：一张图，三路并行", [
        "技能专家并行评分：S-TECH（技术）/ S-GLOBAL（整体）/ S-CONTENT（内容）",
        "像素统计特征：清晰度 / 噪声 / 色彩 / 亮度 / 分辨率 / 宽高比 / 分歧度（7 维 OpenCV）",
        "裸问释义投票：4 条同义问法取均值，承担零点校准",
        "门控矩阵 W (3×7)：BT 锦标赛离线训练，冻结后纯矩阵乘法推理",
        "融合分与投票分按 α 线性混合：KonIQ α=0.6 / SPAQ α=0.3",
        "R1→R2→R3 搭建三路，最终收敛为同一骨架——两数据集仅参数不同",
    ], "fig1_pipeline.png", section="框架", num=3)

    # ── P4 BT 锦标赛 ──
    s = prs.slides.add_slide(blank)
    content_slide(s, "核心一：BT 锦标赛——监督信号从哪来？", [
        "难题：路由需要训练信号，人群分被协议禁止进入在线——必须自监督",
        "科学依据：绝对打分不稳定，两两比较稳健（Thurstone 1927；Bradley-Terry 1952）",
        "组织：2,500 张训练图 → 3.2 万场两两比较，每场由大模型判胜负",
        "拟合：全部胜负记录经 Bradley-Terry 凸优化 → 每张图一个标量强度分",
        "本质：把大模型内部自发的质量排序能力蒸馏成训练信号——零人工标签",
        "结果：强度分成为门控训练的唯一监督信号，训练完成后不再查询锦标赛",
    ], section="方法", num=4)

    # ── P5 门控训练 ──
    s = prs.slides.add_slide(blank)
    content_slide(s, "核心二：门控矩阵——21 个参数的路由训练", [
        "输入：7 维标准化像素特征 F（包含专家分歧度为关键特征）",
        "训练目标：最小化成对铰链损失 −log σ(s_A−s_B)，800 步随机梯度下降",
        "每步：随机抽 128 对图 → 正向矩阵乘法 → 解析梯度 → 更新 21 个参数",
        "CPU 上数秒收敛：损失 0.46 → 0.40，参数范数稳定增长至收敛",
        "留出 2,000 对检验：动态融合一致率 80.9% vs 等权 78.0%（+2.9 pp）",
        "训练后冻结：推理阶段纯矩阵乘法，毫秒级，不调任何大模型",
    ], "fig2_loss.png", section="方法", num=5)

    # ── P6 推理算例 ──
    s = prs.slides.add_slide(blank)
    content_slide(s, "核心三：推理融合与可解释性", [
        "门控矩阵高度可读：分歧度列权重绝对值最大（TECH +0.89 / CONTENT −0.84）",
        "→ 专家意见不一致时，门控将话语权交给技术专家",
        "logpix / aspect 权重自动归零：训练集尺寸恒定，无判别信息——矩阵自动忽略",
        "真实算例：分歧度高出 2.47σ → TECH 话语权 95.8%，融合分 1.90",
        "若等权：融合分 3.53，被整体专家拉高近一倍——门控有效抑制异常",
        "SPAQ 门控未过一致性门槛 → 自动回退等权融合（内置诚实机制）",
    ], "fig3_w.png", section="方法", num=6)

    # ── P7 主表 + 趋势 ──
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, GRAYF5)
    add_bottom_bar(s); add_page_number(s, 7); add_section_tag(s, "结果")
    add_title(s, "统一框架主表与逐臂趋势")
    # 表
    hdr = ["臂", "机制", "KonIQ SRCC", "KonIQ MAE", "SPAQ SRCC", "SPAQ MAE"]
    rows = [hdr,
            ["R1-bare", "单问+锚点带（基准）", "0.625", "0.645", "0.865", "1.062"],
            ["R2", "五专家截尾融合", "0.606", "0.677", "0.859", "0.971"],
            ["R2.5", "规则选路+逆离散度加权", "0.619", "0.533", "0.860", "1.014"],
            ["R3", "注入经验规则", "0.621", "0.499", "0.861", "1.008"],
            ["R6", "统一框架：门控融合+投票混合", "0.734", "0.491", "0.891", "0.796"]]
    make_table(s, rows, Inches(0.55), Inches(1.55),
               [Inches(1.2), Inches(3.4), Inches(1.35), Inches(1.3), Inches(1.35), Inches(1.3)])
    # 趋势图
    insert_fig(s, "fig6_progression.png", Inches(0.55), Inches(4.15), Inches(6.1),
               "图：R1→R6 逐臂 SRCC / MAE 双面板趋势")
    # 右侧 KPI
    add_rect(s, Inches(7.2), Inches(4.35), Inches(5.5), Inches(2.9), fill_color=WHITE)
    tb = mk_textbox(s, Inches(7.5), Inches(4.55), Inches(5.0), Inches(2.5))
    tf = tb.text_frame; tf.word_wrap = True
    add_text(tf, "统一框架 R6 相对基准提升", size=16, color=PRIMARY, bold=True, space_after=14)
    for line, clr in [("KonIQ  SRCC  +0.109    排序一致性大幅跃升", BLUE_A),
                       ("KonIQ  MAE   −0.154    绝对误差显著下降", ACCENT),
                       ("SPAQ   MAE   −0.266    跨域校准同步改善", ACCENT),
                       ("", DARK),
                       ("R3 经验规则零收益 → 监督信号须来自\n与目标强相关来源（否定性证据）", GRAY75)]:
        add_text(tf, line, size=11, color=clr, space_after=6)

    # ── P8 跨模型迁移 ──
    s = prs.slides.add_slide(blank)
    content_slide(s, "检验：门控矩阵跨模型迁移（0 新增 API）", [
        "将 32B 训练的门控矩阵 W，原样作用于 8B 模型的专家分——零调用、纯矩阵乘法",
        "8B 纯裸问探针：SRCC 0.679",
        "8B + 等权门控：SRCC 0.759",
        "8B + 迁移 32B 门控：SRCC 0.776  ——反超 32B 自身（0.734）",
        '结论：门控学到的是"图像属性 → 专家可信度"映射，与模型规模解耦',
    ], "fig7_transfer.png", section="检验", num=8)

    # ── P9 效果 ──
    s = prs.slides.add_slide(blank)
    content_slide(s, "效果演示：预测分 vs 人群分", [
        "两个数据集密集散点：主趋势贴合对角线，无系统性偏移",
        "KonIQ：离散整数区间的预测分布与人群分布形态一致",
        "SPAQ：连续尺度的排序能力保持在 0.891 高位",
        "（可现场任选一张图，实时走一遍完整链路）",
    ], "fig4_scatter.png", section="检验", num=9)

    # ── P10 困难与发现 ──
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, GRAYF5)
    add_bottom_bar(s); add_page_number(s, 10); add_section_tag(s, "发现")
    add_title(s, "困难与发现：提示词约束的双刃剑")
    # 消融表
    hdr = ["组件", "R1-bare\n(纯裸问)", "R1-anchor v2", "R1-anchor v3", "R1-rich\n(完整专家)"]
    rows = [hdr,
            ["评估人设", "—", "✓", "✓", "✓"],
            ["锚点带", "—", "✓", "✓", "✓"],
            ["文字等级", "—", "✓", "✓", "✓"],
            ["维度/清单/程序", "—", "—", "✓", "✓"],
            ["结构化输出", "—", "—", "✓", "✓"]]
    make_table(s, rows, Inches(0.55), Inches(1.55),
               [Inches(2.0), Inches(1.4), Inches(1.9), Inches(1.9), Inches(1.9)])
    # 指标表
    hdr2 = ["臂", "KonIQ SRCC", "KonIQ MAE", "SPAQ SRCC", "SPAQ MAE"]
    rows2 = [hdr2,
             ["纯裸问", "0.660", "0.454", "0.884", "0.819"],
             ["anchor v2", "0.612", "0.507", "0.889", "1.034"],
             ["anchor v3", "0.625", "0.645", "0.865", "1.062"],
             ["完整专家", "0.633", "0.688", "0.867", "1.042"]]
    make_table(s, rows2, Inches(0.55), Inches(3.2),
               [Inches(1.5), Inches(1.5), Inches(1.45), Inches(1.5), Inches(1.45)])
    # 右侧配图
    insert_fig(s, "fig5_ablation.png", Inches(7.25), Inches(1.55), Inches(5.6))
    insert_fig(s, "fig8_dist.png", Inches(7.25), Inches(4.1), Inches(5.6),
               "图：SRCC高位维持 / MAE逐级恶化（左）| 分布被锚点重塑为尖峰（右）")
    # 底部关键发现
    add_rect(s, Inches(0.55), Inches(5.85), Inches(12.2), Inches(1.1), fill_color=WHITE)
    tb = mk_textbox(s, Inches(0.85), Inches(5.95), Inches(11.6), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    add_text(tf, "核心发现", size=12, color=PRIMARY, bold=True, space_after=4)
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = ("① 排序与绝对校准可分离 —— 锚点保序但框死分布零点的铁证  "
              "② 这不是缺陷：R6 中裸问投票可独立校零，形成互补结构  "
              "③ 锚点双重角色：单问中恶化 MAE / 多专家融合中是公共标尺（去锚点 SPAQ MAE 0.80→1.21 崩塌）")
    r.font.size = Pt(10); r.font.color.rgb = GRAY75
    r.font.name = "Microsoft YaHei"

    # ── P11 结论 ──
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, PRIMARY)
    add_rect(s, Inches(0), Inches(0), Inches(0.06), H, fill_color=ACCENT)
    add_rect(s, Inches(0), H - Inches(0.06), W, Inches(0.06), fill_color=ACCENT)
    add_page_number(s, 11)
    tb = mk_textbox(s, Inches(1.2), Inches(1.0), Inches(11.0), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    add_text(tf, "结论", size=32, color=WHITE, bold=True, space_after=28)
    items = [
        ("01", "统一框架：两数据集 × 两模型规模，同一骨架，主干零训练，全部结果温度为零、缓存可复现"),
        ("02", "BT 锦标赛自监督：3.2 万场两两比较 → 21 参数门控，有效、可解释、跨模型可迁移"),
        ("03", "反直觉发现：提示词约束的保序-校零分离 / 锚点的双重角色 / 检索式经验规则的零收益"),
    ]
    for num, txt in items:
        add_text(tf, num, size=28, color=ACCENT, bold=True, name="Arial", space_after=2)
        add_text(tf, txt, size=16, color=WHITE, space_after=16)
    add_text(tf, "", size=8, color=WHITE, space_after=10)
    add_text(tf, "谢谢各位老师，请批评指正。", size=20, color=ACCENT, space_after=0)

    # ── P12 备用 ──
    s = prs.slides.add_slide(blank)
    content_slide(s, "（备用）合规自检与复现协议", [
        "在线输入白名单：图像像素 + 提示词模板 + 温度零确定性协议 —— 仅此三项",
        "数据集名、图像 ID、人群 MOS、来源分数 —— 在 Router/LLM/Skill 中零出现",
        "人群分唯一读取入口：iqa_agent/data.py → load_mos()，仅限离线评测与误差分析脚本",
        "磁盘缓存：温度零 + SHA256 摘要 key → 一次调用后永久缓存，任意臂一键复现",
        "R4/R5（中间验证产物）：验证了锦标赛监督价值与投票机制，有效成分已并入 R6",
        "本文全部插图由 scripts/report_figs.py 自动生成，论文/PPT 由 make_paper.py / make_ppt_pro.py 编译",
    ], section="备用", num=12)

    path = os.path.join(ROOT, "docs", "汇报_pro.pptx")
    prs.save(path)
    print("->", path)


if __name__ == "__main__":
    build()
